from pptx import Presentation
from pptx.util import Pt

# Erstellen der Präsentation
prs = Presentation()

# Folieninhalte definieren
content = [
    {
        "title": "Frühes 20. Jahrhundert",
        "points": [
            "Neoklassische Theorie: Fokus auf Marktgleichgewicht und Preisbildung durch Angebot und Nachfrage.",
            "Marxistische Wirtschaftstheorie: Analyse des Kapitalismus und der Einkommensungleichheit.",
            "Positivismus: Anwendung naturwissenschaftlicher Methoden auf soziale Phänomene.",
            "Soziologie von Émile Durkheim: Konzepte wie 'kollektives Bewusstsein' und soziale Normen."
        ]
    },
    {
        "title": "Mitte bis Ende des 20. Jahrhunderts",
        "points": [
            "Keynesianismus: Rolle des Staates in der Wirtschaft und staatliche Investitionen.",
            "Monetarismus: Bedeutung der Geldmenge für Inflation und wirtschaftliche Stabilität.",
            "Strukturfunktionalismus: Gesellschaft als System interdependenter Teile.",
            "Kritische Theorie: Analyse von Machtstrukturen und kultureller Hegemonie."
        ]
    },
    {
        "title": "Spätes 20. Jahrhundert bis heute",
        "points": [
            "Neoliberalismus: Deregulierung, Privatisierung und freier Handel als Wachstumsförderer.",
            "Verhaltensökonomie: Einfluss von psychologischen Faktoren auf wirtschaftliche Entscheidungen.",
            "Poststrukturalismus: Dekonstruktion sozialer Normen und Analyse von Sprache und Macht.",
            "Globalisierungstheorien: Auswirkungen der Globalisierung auf nationale Souveränität und soziale Gerechtigkeit."
        ]
    }
]

# Folien hinzufügen
for slide_content in content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = slide_content["title"]
    body = slide.shapes.placeholders[1]

    for point in slide_content["points"]:
        p = body.text_frame.add_paragraph()
        p.text = point
        p.level = 0  # Setze die Ebene auf 0 für den Haupttext
        p.font.size = Pt(18)

# Präsentation speichern
pptx_file = "C:\WiSo_Praesentation.pptx"
prs.save(pptx_file)
